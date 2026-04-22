"""
FastAPI Backend Server for Financial Analysis Agents
Provides REST API and Server-Sent Events (SSE) for streaming responses
"""
from __future__ import annotations
import os
import json
import asyncio
import logging
import threading
import uuid as uuid_mod
from collections import OrderedDict
from typing import Optional, AsyncGenerator, Any, Dict, List
from datetime import datetime, timedelta, timezone
from contextlib import asynccontextmanager
from fastapi import FastAPI, HTTPException, Request, UploadFile, File, Depends
from fastapi.concurrency import run_in_threadpool
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, JSONResponse, Response, FileResponse
from fastapi.staticfiles import StaticFiles
from fastapi.exceptions import RequestValidationError
from pydantic import BaseModel, ValidationError
from dotenv import load_dotenv
import sys
import requests
import re
import secrets
from sqlalchemy import select, or_
from sqlalchemy.ext.asyncio import AsyncSession

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Add parent directory to path to import agents
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from backend.config import (
    SSE_CHUNK_SIZE, SSE_STREAM_DELAY_SECONDS,
    CHART_PERIOD_DAYS, CORS_ORIGINS,
)
from shared.ticker_utils import extract_ticker as _extract_ticker_shared
from backend.callbacks.streaming import StreamingCallbackHandler
from backend.database import init_db, get_db, SyncSessionLocal
from backend.models import Session as DBSession, DBMessage, Analysis, Watchlist, WatchlistTicker, Project, ProjectSession, ProjectDocument
from backend.project_config import normalize_project_config
from agents.finance_qa_agent import create_finance_qa_agent
from agents.market_agent import create_market_agent
from agents.portfolio_agent import create_portfolio_agent
from agents.earnings_agent import create_earnings_agent
from arena.arena_agent import ArenaAgent
from arena.run import run_arena
from arena.progress import set_arena_queue, clear_arena_queue
from arena.output import extract_structured_memo

# Load environment variables from parent directory
load_dotenv(os.path.join(os.path.dirname(__file__), '..', '.env'))

@asynccontextmanager
async def lifespan(app: FastAPI):
    """Initialize resources on startup."""
    await init_db()
    logger.info("Database initialized")
    # Pre-load Chroma embedding model so the ~90MB download happens before first request
    try:
        from data.chroma_client import ProjectChromaClient
        _ = ProjectChromaClient()
        logger.info("ProjectChromaClient initialised")
    except Exception as _e:
        logger.warning(f"ProjectChromaClient pre-load failed (non-fatal): {_e}")
    # Start heartbeat scheduler
    try:
        from backend.scheduler import start_scheduler, stop_scheduler
        await start_scheduler()
        logger.info("Heartbeat scheduler started")
        yield
        await stop_scheduler()
    except Exception as _e:
        logger.warning(f"Scheduler startup failed (non-fatal): {_e}")
        yield


# Rate limiting — in-memory, single-worker only (Phase 1).
# Phase 2: replace with Redis INCR + TTL (see TODOS.md).
from collections import defaultdict as _defaultdict

_memo_ip_requests: dict = _defaultdict(list)
_IP_MEMO_LIMIT = 5
_IP_WINDOW_SECONDS = 3600  # 1 hour

_memo_daily_count = 0
_memo_daily_reset_date: Optional[str] = None
MEMO_DAILY_CAP = 50


def _check_memo_rate_limits(client_ip: str) -> None:
    """Check both per-IP (5/hour) and global daily (50/day) caps. Raises HTTP 429 on breach."""
    global _memo_daily_count, _memo_daily_reset_date

    # Per-IP sliding window
    now = datetime.now(timezone.utc).timestamp()
    window_start = now - _IP_WINDOW_SECONDS
    _memo_ip_requests[client_ip] = [t for t in _memo_ip_requests[client_ip] if t > window_start]
    if len(_memo_ip_requests[client_ip]) >= _IP_MEMO_LIMIT:
        raise HTTPException(status_code=429, detail="Rate limit reached — try again later")
    _memo_ip_requests[client_ip].append(now)

    # Global daily cap
    today = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    if _memo_daily_reset_date != today:
        _memo_daily_count = 0
        _memo_daily_reset_date = today
    _memo_daily_count += 1
    if _memo_daily_count > MEMO_DAILY_CAP:
        raise HTTPException(status_code=429, detail="Rate limit reached — try again later")


# Path to frontend build — defined early so root route can reference it
FRONTEND_BUILD_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "frontend", "dist")

# Initialize FastAPI app
app = FastAPI(title="Financial Analysis API", version="1.0.0", lifespan=lifespan)

# Register scheduled agents router
from backend.scheduled_agents_router import router as scheduled_agents_router
app.include_router(scheduled_agents_router)

# Register CIO router
from backend.cio_router import router as cio_router
app.include_router(cio_router)



# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=CORS_ORIGINS,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Add validation error handler for better debugging
@app.exception_handler(RequestValidationError)
async def validation_exception_handler(request: Request, exc: RequestValidationError):
    """Log validation errors for debugging"""
    body = await request.body()
    logger.error(f"Validation error on {request.url.path}")
    logger.error(f"Request body: {body}")
    logger.error(f"Validation errors: {exc.errors()}")
    return JSONResponse(
        status_code=422,
        content={"detail": exc.errors(), "body": body.decode('utf-8')}
    )

# Store active agents with bounded LRU eviction
_AGENTS_CACHE_MAX = 100
agents_cache: OrderedDict = OrderedDict()
agents_cache_lock = threading.Lock()
SESSION_SCOPED_AGENT_TYPES = frozenset({"research", "market", "earnings", "arena"})

# Hold strong references to fire-and-forget tasks so they aren't GC'd before completion
_background_tasks: set = set()


def _fire_and_forget(coro) -> asyncio.Task:
    """Schedule a coroutine as a background task, keeping a reference until it completes."""
    task = asyncio.create_task(coro)
    _background_tasks.add(task)
    task.add_done_callback(_background_tasks.discard)
    return task

# Map agent types to their fallback methods (when agent_executor is not available)
AGENT_FALLBACK_METHODS = {
    "research": "chat",  # Research uses 'chat' instead of 'analyze'
    "market": "analyze",
    "portfolio": "analyze",
    "earnings": "analyze",
    "arena": "analyze",
}


def extract_ticker_from_query(query: str, is_followup: bool = False) -> Optional[str]:
    """Delegate to shared.ticker_utils.extract_ticker."""
    return _extract_ticker_shared(query, is_followup=is_followup)


class ChatMessage(BaseModel):
    """Chat message model"""
    message: str
    agent_type: str = "research"  # research, market, portfolio, earnings, arena
    model: str = "claude-sonnet-4-6"
    session_id: Optional[str] = None
    is_followup: bool = False
    # Persistence metadata (populated by frontend after receiving session_id)
    persist: bool = True   # set False to skip DB write (e.g. health checks)
    project_id: Optional[str] = None  # Set when chat is inside a project workspace


class ChatResponse(BaseModel):
    """Chat response model"""
    response: str
    agent_type: str
    timestamp: str
    session_id: str


def _build_agent_cache_key(agent_type: str, model: str, session_id: Optional[str]) -> Optional[str]:
    """Use session-scoped cache keys for stateful agents to prevent context leakage."""
    if agent_type in SESSION_SCOPED_AGENT_TYPES:
        if not session_id:
            return None
        return f"{agent_type}_{model}_{session_id}"
    return f"{agent_type}_{model}"


def _create_agent_instance(agent_type: str, model: str):
    """Create a single agent instance for the requested type."""
    if agent_type == "research":
        return create_finance_qa_agent(model=model, db_session_factory=SyncSessionLocal)
    if agent_type == "market":
        return create_market_agent(model=model)
    if agent_type == "portfolio":
        return create_portfolio_agent(model=model)
    if agent_type == "earnings":
        return create_earnings_agent(model=model)
    if agent_type == "arena":
        return ArenaAgent()
    raise ValueError(f"Unknown agent type: {agent_type}")


def get_or_create_agent(agent_type: str, model: str, session_id: Optional[str] = None):
    """Get cached agent or create a fresh instance when session scoping is required."""
    cache_key = _build_agent_cache_key(agent_type, model, session_id)

    try:
        if cache_key is None:
            return _create_agent_instance(agent_type, model)

        with agents_cache_lock:
            if cache_key in agents_cache:
                agents_cache.move_to_end(cache_key)  # LRU: mark as recently used
                return agents_cache[cache_key]
            agent = _create_agent_instance(agent_type, model)
            agents_cache[cache_key] = agent
            while len(agents_cache) > _AGENTS_CACHE_MAX:
                agents_cache.popitem(last=False)  # evict oldest (LRU)
            return agent
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to create agent: {str(e)}")


def _ensure_str_response(value) -> str:
    """Normalize Anthropic content blocks (list) or other types to a plain string.

    ChatAnthropic returns AIMessage.content as a list of content blocks
    (e.g., [{"type": "text", "text": "..."}]) instead of a plain string.
    This function extracts the text from all blocks and returns a single string.
    """
    if isinstance(value, str):
        return value
    if isinstance(value, list):
        parts = []
        for block in value:
            if isinstance(block, dict):
                parts.append(block.get("text", ""))
            elif isinstance(block, str):
                parts.append(block)
            else:
                text = getattr(block, "text", None)
                if text:
                    parts.append(str(text))
        return "".join(parts)
    return str(value) if value is not None else ""


async def run_agent_with_callbacks(agent, message: str, agent_type: str, queue: asyncio.Queue, is_followup: bool = False):
    """
    Run agent in executor with callback handler.

    Executes the agent synchronously in a thread pool executor while streaming
    events (thoughts, tool calls, results) to an async queue for SSE delivery.

    Args:
        agent: LangChain agent instance (with agent_executor or fallback method)
        message: User's input message to process
        agent_type: One of 'research', 'market', 'portfolio', 'earnings', 'arena'
        queue: Async queue for streaming events to SSE response
        is_followup: Whether this is a follow-up question (earnings agent only)
    """
    loop = asyncio.get_running_loop()
    callback = StreamingCallbackHandler(queue)

    try:
        # Validate agent type
        if agent_type not in AGENT_FALLBACK_METHODS:
            raise ValueError(f"Unknown agent type: {agent_type}")

        # Inject progress queue for agents that use direct tool calls (bypassing LangChain callbacks)
        if agent_type in ("earnings", "graph", "arena"):
            agent._progress_queue = queue
            agent._progress_loop = loop

        # Use _invoke() if available — it resets per-request callback state before calling
        # agent_executor, ensuring step counters don't accumulate on cached agent instances.
        # Fall back to agent_executor.invoke() for agents that don't implement _invoke().
        if hasattr(agent, '_invoke'):
            input_dict = {"input": message}
            if is_followup and agent_type == "earnings":
                input_dict["followup"] = True
            response = await loop.run_in_executor(
                None,
                lambda: agent._invoke(input_dict, [callback])
            )
        elif hasattr(agent, 'agent_executor'):
            input_dict = {"input": message}
            if is_followup and agent_type == "earnings":
                input_dict["followup"] = True
            response = await loop.run_in_executor(
                None,
                lambda: agent.agent_executor.invoke(
                    input_dict,
                    config={"callbacks": [callback]}
                )["output"]
            )
        else:
            # Fallback to agent's direct method (analyze or chat)
            fallback_method_name = AGENT_FALLBACK_METHODS[agent_type]
            fallback_method = getattr(agent, fallback_method_name)
            response = await loop.run_in_executor(None, fallback_method, message)

        # Clean up progress queue
        if agent_type in ("earnings", "graph", "arena"):
            agent._progress_queue = None
            agent._progress_loop = None

        # Normalize response to string — Anthropic returns list of content blocks
        response = _ensure_str_response(response)

        await queue.put({"type": "response", "content": response})
        await queue.put({"type": "done"})

    except Exception as e:
        # Clean up progress queue on error too
        if agent_type in ("earnings", "graph", "arena"):
            agent._progress_queue = None
            agent._progress_loop = None
        await queue.put({"type": "error", "error": str(e)})


async def run_project_graph_with_callbacks(
    adapter,
    input_dict: dict,
    queue: asyncio.Queue,
    state_container: dict,
) -> None:
    """Run ProjectAnalysisGraph in executor, stream progress events via queue.

    Final response is placed in queue as {"type": "response", "content": ...}.
    Extracted state (including memory_patch) is written to state_container["state"].
    """
    loop = asyncio.get_running_loop()
    graph_instance = adapter.graph_instance
    graph_instance._progress_queue = queue
    graph_instance._progress_loop = loop

    try:
        result = await loop.run_in_executor(None, lambda: adapter.invoke(input_dict))
        state_container["state"] = result.get("_state", {}) if result else {}
        output = result.get("output", "") if result else ""
        await queue.put({"type": "response", "content": output})
        await queue.put({"type": "done"})
    except Exception as exc:
        state_container["state"] = {}
        await queue.put({"type": "error", "error": str(exc)})
    finally:
        graph_instance._progress_queue = None
        graph_instance._progress_loop = None


async def route_agent_for_message(message: str) -> str:
    """Use Claude Haiku to classify a user message to the best chat agent.

    Returns one of: 'research', 'analyst', 'market'
    Falls back to 'research' on any error.
    """
    import anthropic
    client = anthropic.AsyncAnthropic()

    routing_prompt = (
        "Route this financial query to the best agent. Reply with ONLY one word.\n\n"
        "Agents:\n"
        "- earnings: Earnings reports, quarterly results, EPS beats/misses, earnings call analysis, "
        "analyst estimates, earnings surprises, revenue guidance, management commentary from calls. "
        "ONLY use this if the query explicitly mentions earnings, EPS, quarterly results, or beats/misses "
        "AND names a specific company or ticker.\n"
        "- analyst: Deep equity analysis, investment thesis, moat/competitive analysis, "
        "buy/sell recommendation, 'should I invest' questions, stock screeners\n"
        "- market: Market conditions, S&P 500/NASDAQ/Dow indices, VIX, sector rotation, "
        "macro trends, Fed policy, inflation, recession risk, market sentiment\n"
        "- research: Everything else — company info, stock comparisons, financial metrics, "
        "quick Q&A, revenue/profit data, R&D spending, product roadmaps, follow-up questions "
        "that reference earlier context without naming a new company, and any question that does "
        "NOT explicitly name a specific company or ticker.\n\n"
        "IMPORTANT: If the query looks like a follow-up (uses words like 'these', 'this', 'their', "
        "'the company', 'it', 'they' without naming a company), route to 'research'.\n\n"
        f"Message: {message}\n\n"
        "Agent:"
    )

    try:
        response = await client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=10,
            messages=[{"role": "user", "content": routing_prompt}],
        )
        agent = _ensure_str_response(response.content).strip().lower()
        return agent if agent in ("research", "analyst", "market", "earnings") else "research"
    except Exception as e:
        logger.warning(f"Auto-routing failed, defaulting to research: {e}")
        return "research"


def _requests_get_json(url: str, *, params: Dict[str, Any], timeout: int = 10) -> Any:
    """Synchronous helper for requests-based APIs."""
    response = requests.get(url, params=params, timeout=timeout)
    response.raise_for_status()
    return response.json()


async def _fetch_json(url: str, *, params: Dict[str, Any], timeout: int = 10) -> Any:
    """Run blocking HTTP requests in a threadpool so async endpoints stay responsive."""
    return await run_in_threadpool(_requests_get_json, url, params=params, timeout=timeout)


async def generate_follow_up_questions(message: str, response: str, agent_type: str) -> list[str]:
    """Generate 3 contextual follow-up questions using Claude Haiku (fast + cheap).

    Returns an empty list if generation fails (non-blocking).
    """
    import anthropic

    client = anthropic.AsyncAnthropic()

    agent_context = {
        "analyst": "comprehensive equity research",
        "graph": "structured equity research",
        "research": "financial research",
        "market": "market analysis and sentiment",
        "portfolio": "portfolio analysis and optimization",
        "earnings": "earnings analysis and quarterly trends",
    }.get(agent_type, "financial analysis")

    try:
        result = await client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=256,
            system=(
                f"You are a follow-up question generator for a {agent_context} tool. "
                "Generate exactly 3 brief follow-up questions an investor might ask next, "
                "based on the conversation. Each question should explore a different angle. "
                "Return only the questions, one per line, no numbering or bullets."
            ),
            messages=[
                {
                    "role": "user",
                    "content": f"User asked: {message}\n\nAssistant responded (excerpt):\n{response[:2000]}",
                }
            ],
        )
        if not result.content:
            return []
        text = result.content[0].text.strip()
    except Exception:
        return []

    questions = [q.strip() for q in text.split("\n") if q.strip()]
    return questions[:3]


async def stream_agent_response(
    message: str,
    agent_type: str,
    model: str,
    is_followup: bool = False,
    session_id: Optional[str] = None,
    persist: bool = True,
    project_id: Optional[str] = None,
) -> AsyncGenerator[str, None]:
    """Stream agent response using Server-Sent Events with thinking process"""
    queue = asyncio.Queue()

    # ── Project session path ──────────────────────────────────────────────────
    if project_id:
        try:
            from backend.database import AsyncSessionLocal
            from backend.context_assembly import assemble_project_context
            from backend.project_router import route_for_project
            from agents.project_agent import ProjectAnalysisGraph, ProjectAnalysisGraphAdapter
            from data.chroma_client import ProjectChromaClient

            # 1. Assemble context and load project config
            async with AsyncSessionLocal() as _db:
                _chroma = ProjectChromaClient()
                context_block = await assemble_project_context(project_id, message, _db, _chroma)
                _proj_result = await _db.execute(select(Project).where(Project.id == project_id))
                _proj = _proj_result.scalar_one_or_none()
                if not _proj:
                    yield f"data: {json.dumps({'type': 'error', 'error': 'Project not found'})}\n\n"
                    return
                project_config: dict = json.loads(_proj.config) if _proj.config else {}

            # 2. Route query to agents
            routing_decision_obj = await route_for_project(message, context_block, project_config)
            routing_decision = {
                "agents": [
                    {"agent_type": a.agent_type, "task": a.task}
                    for a in routing_decision_obj.agents
                ],
                "reasoning": routing_decision_obj.reasoning,
            }

            yield f"data: {json.dumps({'type': 'routing_decision', 'agent': 'project', 'routing': routing_decision})}\n\n"
            yield f"data: {json.dumps({'type': 'start', 'agent': 'project'})}\n\n"

            # 3. Run graph
            graph = ProjectAnalysisGraph()
            adapter = ProjectAnalysisGraphAdapter(graph)
            state_container: dict = {}
            task = asyncio.create_task(
                run_project_graph_with_callbacks(
                    adapter,
                    {
                        "input": message,
                        "project_id": project_id,
                        "context_block": context_block,
                        "routing_decision": routing_decision,
                    },
                    queue,
                    state_container,
                )
            )

            # 4. Drain queue (same pattern as regular path)
            collected_response = ""
            while True:
                event = await queue.get()
                if event["type"] == "done":
                    break
                elif event["type"] == "error":
                    yield f"data: {json.dumps({'type': 'error', 'error': event['error']})}\n\n"
                    task.cancel()
                    break
                elif event["type"] == "response":
                    collected_response = event["content"]
                    for i in range(0, len(collected_response), SSE_CHUNK_SIZE):
                        chunk = collected_response[i:i + SSE_CHUNK_SIZE]
                        yield f"data: {json.dumps({'type': 'content', 'content': chunk})}\n\n"
                        await asyncio.sleep(SSE_STREAM_DELAY_SECONDS)
                else:
                    yield f"data: {json.dumps(event)}\n\n"

            try:
                await task
            except asyncio.CancelledError:
                pass

            # 5a. Non-blocking memory update background task
            final_state = state_container.get("state", {})
            memory_patch = final_state.get("memory_patch") or {}
            if memory_patch:
                try:
                    _fire_and_forget(_run_memory_update(project_id, memory_patch))
                except Exception as _me:
                    logger.warning("Failed to schedule memory update for project %s: %s", project_id, _me)

            # 5. Persist with project_id linkage
            ticker = extract_ticker_from_query(message, is_followup=is_followup)
            if persist and collected_response:
                try:
                    _fire_and_forget(
                        _persist_conversation(
                            session_id=session_id,
                            user_message=message,
                            assistant_response=collected_response,
                            agent_type="project",
                            ticker=ticker,
                            thinking_steps=[],
                            follow_ups=[],
                            project_id=project_id,
                        )
                    )
                except Exception as _pe:
                    logger.warning(f"Failed to schedule project persistence: {_pe}")

            yield f"data: {json.dumps({'type': 'end'})}\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'type': 'error', 'error': f'Project analysis error: {str(e)}'})}\n\n"
        return
    # ── End project path ──────────────────────────────────────────────────────

    try:
        # Auto-route when agent_type is "auto" — use Claude Haiku to classify the query
        resolved_agent_type = agent_type
        if agent_type == "auto":
            resolved_agent_type = await route_agent_for_message(message)
            logger.info(f"[AUTO-ROUTE] '{message[:60]}...' → {resolved_agent_type}")
            yield f"data: {json.dumps({'type': 'routing_decision', 'agent': resolved_agent_type})}\n\n"

        agent = get_or_create_agent(resolved_agent_type, model, session_id=session_id)

        # Send start event
        yield f"data: {json.dumps({'type': 'start', 'agent': resolved_agent_type})}\n\n"

        # Extract ticker from user query and send as metadata
        ticker = extract_ticker_from_query(message, is_followup=is_followup)
        if ticker:
            print(f"[INFO] Detected ticker from query: {ticker}")
            yield f"data: {json.dumps({'type': 'ticker_metadata', 'ticker': ticker})}\n\n"

        # Start agent execution in background
        task = asyncio.create_task(run_agent_with_callbacks(agent, message, resolved_agent_type, queue, is_followup))

        # Accumulate response text for follow-up generation
        collected_response = ""
        collected_thinking: list = []
        collected_charts: dict = {}

        # Stream events from queue
        while True:
            event = await queue.get()

            if event["type"] == "done":
                break
            elif event["type"] == "error":
                yield f"data: {json.dumps({'type': 'error', 'error': event['error']})}\n\n"
                break
            elif event["type"] == "response":
                # Stream the final response in chunks
                response = event["content"]
                collected_response = response
                for i in range(0, len(response), SSE_CHUNK_SIZE):
                    chunk = response[i:i + SSE_CHUNK_SIZE]
                    yield f"data: {json.dumps({'type': 'content', 'content': chunk})}\n\n"
                    await asyncio.sleep(SSE_STREAM_DELAY_SECONDS)
            elif event["type"] == "chart_data":
                chart_id = event.get("id")
                if chart_id:
                    collected_charts[chart_id] = event
                yield f"data: {json.dumps(event)}\n\n"
            else:
                # Collect thinking steps for persistence
                if event["type"] in ("thought", "tool", "tool_result"):
                    collected_thinking.append(event)
                # Stream thinking events (thought, tool, tool_result)
                yield f"data: {json.dumps(event)}\n\n"

        # Wait for task to complete
        await task

        # Generate follow-up questions (non-blocking — failures are silently ignored)
        follow_ups: list[str] = []
        if collected_response:
            try:
                follow_ups = await generate_follow_up_questions(message, collected_response, resolved_agent_type)
                if follow_ups:
                    yield f"data: {json.dumps({'type': 'follow_ups', 'questions': follow_ups})}\n\n"
            except Exception:
                pass

        # Persist session + messages + optional analysis to DB
        if persist and collected_response:
            try:
                _fire_and_forget(
                    _persist_conversation(
                        session_id=session_id,
                        user_message=message,
                        assistant_response=collected_response,
                        agent_type=resolved_agent_type,
                        ticker=ticker,
                        thinking_steps=collected_thinking,
                        follow_ups=follow_ups,
                        chart_specs=collected_charts,
                    )
                )
            except Exception as e:
                logger.warning(f"Failed to schedule persistence task: {e}")

        # Send end event
        yield f"data: {json.dumps({'type': 'end'})}\n\n"

    except Exception as e:
        error_msg = f"Error: {str(e)}"
        yield f"data: {json.dumps({'type': 'error', 'error': error_msg})}\n\n"


# Agent types that should auto-save to the analyses library
_ANALYSIS_AGENT_TYPES = {"earnings"}


async def _run_memory_update(project_id: str, memory_patch: dict) -> None:
    """Background wrapper: open own DB session and apply memory patch.

    Errors are logged but never raised — memory update failure must not affect
    the user-visible response.
    """
    try:
        from backend.database import AsyncSessionLocal
        from data.project_memory import update_project_memory
        async with AsyncSessionLocal() as db:
            await update_project_memory(project_id, memory_patch, db)
        logger.info("Memory update completed for project %s", project_id)
    except Exception as exc:  # noqa: BLE001
        logger.error("_run_memory_update failed for project %s: %s", project_id, exc, exc_info=True)


async def _persist_conversation(
    session_id: Optional[str],
    user_message: str,
    assistant_response: str,
    agent_type: str,
    ticker: Optional[str],
    thinking_steps: list,
    follow_ups: list[str],
    chart_specs: Optional[dict] = None,
    project_id: Optional[str] = None,
) -> None:
    """Persist session, messages, and optional analysis to the database."""
    from backend.database import AsyncSessionLocal
    import uuid as _uuid_mod

    try:
        async with AsyncSessionLocal() as db:
            # Upsert session
            sid = session_id or str(_uuid_mod.uuid4())
            result = await db.execute(select(DBSession).where(DBSession.id == sid))
            db_session = result.scalar_one_or_none()

            if db_session is None:
                title = user_message[:60].strip() or "Conversation"
                db_session = DBSession(
                    id=sid,
                    title=title,
                    agent_type=agent_type,
                )
                db.add(db_session)
            else:
                db_session.last_active_at = datetime.now(timezone.utc)
                db_session.agent_type = agent_type

            # Insert user message
            user_msg = DBMessage(
                session_id=sid,
                role="user",
                content=user_message,
                agent_type=agent_type,
                ticker=ticker,
            )
            db.add(user_msg)
            await db.flush()  # get user_msg.id

            # Insert assistant message
            assistant_msg = DBMessage(
                session_id=sid,
                role="assistant",
                content=assistant_response,
                agent_type=agent_type,
                ticker=ticker,
                thinking_steps=json.dumps(thinking_steps) if thinking_steps else None,
                follow_ups=json.dumps(follow_ups) if follow_ups else None,
                chart_specs=json.dumps(chart_specs) if chart_specs else None,
            )
            db.add(assistant_msg)
            await db.flush()

            # Auto-save analysis for qualifying agent types
            if agent_type in _ANALYSIS_AGENT_TYPES and ticker:
                month_str = datetime.now(timezone.utc).strftime("%b %Y")
                agent_label = {
                    "analyst": "Equity Analyst",
                    "earnings": "Earnings",
                    "graph": "Equity Research",
                }.get(agent_type, agent_type.title())
                analysis = Analysis(
                    session_id=sid,
                    message_id=assistant_msg.id,
                    ticker=ticker,
                    agent_type=agent_type,
                    title=f"{ticker} {agent_label} Analysis — {month_str}",
                    content=assistant_response,
                    tags="[]",
                )
                db.add(analysis)

            # Link session to project (upsert — UNIQUE constraint prevents duplicates)
            if project_id:
                from sqlalchemy.dialects.sqlite import insert as sqlite_insert
                link_stmt = sqlite_insert(ProjectSession).values(
                    id=str(_uuid_mod.uuid4()),
                    project_id=project_id,
                    session_id=sid,
                    created_at=datetime.now(timezone.utc),
                ).on_conflict_do_nothing()
                await db.execute(link_stmt)

            await db.commit()
            logger.info(f"[DB] Persisted session {sid} with {agent_type} message")
    except Exception as e:
        logger.error(f"[DB] Persistence error: {e}")


@app.get("/api/health")
async def health():
    """Health check endpoint"""
    return {
        "status": "online",
        "service": "Financial Analysis API",
        "version": "1.0.0",
        "agents": ["research", "market", "portfolio", "earnings", "arena"]
    }

@app.get("/")
async def root():
    """Serve frontend or health check"""
    index_path = os.path.join(FRONTEND_BUILD_DIR, "index.html")
    if os.path.isfile(index_path):
        return FileResponse(index_path)
    return {
        "status": "online",
        "service": "Financial Analysis API",
        "version": "1.0.0",
        "agents": ["research", "market", "portfolio", "earnings", "arena"]
    }


@app.get("/agents")
async def list_agents():
    """List available agents"""
    return {
        "agents": [
            {
                "id": "research",
                "name": "Finance Q&A",
                "description": "Your personal equity analyst. Ask about any stock — get a quick brief, then dig into valuation, earnings, competitive position, SEC filings, or management commentary.",
                "example": "I'm looking at NVDA — what should I know?"
            },
            {
                "id": "market",
                "name": "Market Analyst",
                "description": "Market conditions, sentiment, and sector analysis",
                "example": "What's the current market sentiment?"
            },
            {
                "id": "portfolio",
                "name": "Portfolio Analyzer",
                "description": "Portfolio analysis with metrics, diversification, and tax optimization",
                "example": "Analyze my portfolio: [{'ticker': 'AAPL', 'shares': 100, 'cost_basis': 150.00}, {'ticker': 'MSFT', 'shares': 50, 'cost_basis': 250.00}]"
            },
            {
                "id": "earnings",
                "name": "Earnings Analyst",
                "description": "Fast earnings-focused equity research (15 min) with quarterly trends and estimates",
                "example": "Analyze NVDA's latest earnings and forward outlook"
            },
            {
                "id": "arena",
                "name": "Investment Committee",
                "description": "Multi-agent debate: Fundamental, Risk, Quant, Macro, and Sentiment analysts debate to conviction-rated investment memo",
                "example": "Should we long NVDA? Full IC review."
            }
        ]
    }


@app.post("/chat/stream")
async def chat_stream(chat_message: ChatMessage):
    """Stream chat response using Server-Sent Events"""
    logger.info(f"[CHAT_STREAM] Received request - agent_type: {chat_message.agent_type}, model: {chat_message.model}, message length: {len(chat_message.message)}")
    return StreamingResponse(
        stream_agent_response(
            chat_message.message,
            chat_message.agent_type,
            chat_message.model,
            chat_message.is_followup,
            session_id=chat_message.session_id,
            persist=chat_message.persist,
            project_id=chat_message.project_id,
        ),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no"
        }
    )


@app.post("/chat", response_model=ChatResponse)
async def chat(chat_message: ChatMessage):
    """Non-streaming chat endpoint (for simple requests)"""
    try:
        session_id = chat_message.session_id
        if chat_message.agent_type in SESSION_SCOPED_AGENT_TYPES and not session_id:
            session_id = str(uuid_mod.uuid4())

        agent = get_or_create_agent(chat_message.agent_type, chat_message.model, session_id=session_id)

        # Get response synchronously.
        # Use _invoke() where available (resets per-request state, no CLI stdout callbacks).
        # Research agent uses its own 'chat' method; other agents without _invoke fall back to analyze().
        if chat_message.agent_type == "research":
            response = await run_in_threadpool(agent.chat, chat_message.message)
        elif hasattr(agent, '_invoke'):
            response = await run_in_threadpool(
                lambda: agent._invoke({"input": chat_message.message}, [])
            )
        else:
            response = await run_in_threadpool(agent.analyze, chat_message.message)

        response = _ensure_str_response(response)

        return ChatResponse(
            response=response,
            agent_type=chat_message.agent_type,
            timestamp=datetime.now().isoformat(),
            session_id=session_id or "default"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


def derive_verdict(
    consensus_score: float,
    agent_signals: dict,
    next_action: str = "",
) -> tuple:
    """
    Derive a BUY / WATCH / PASS verdict from the final ThesisState.
    Returns (verdict: str, confidence: float).
    """
    bullish_count = sum(
        1 for s in agent_signals.values() if s.get("view") == "bullish"
    )
    bearish_count = sum(
        1 for s in agent_signals.values() if s.get("view") == "bearish"
    )

    if next_action == "escalate_to_human":
        return ("WATCH", consensus_score)

    if consensus_score >= 0.70 and bullish_count >= 3:
        verdict = "BUY"
    elif consensus_score <= 0.40 or bearish_count >= 3:
        verdict = "PASS"
    else:
        verdict = "WATCH"

    return (verdict, consensus_score)


class MemoRequest(BaseModel):
    ticker: str
    query_mode: str = "full_ic"


def _run_arena_in_worker(
    query: str,
    ticker: str,
    query_mode: str,
    queue: asyncio.Queue,
    loop: asyncio.AbstractEventLoop,
) -> dict:
    """Thin wrapper executed inside the threadpool worker.
    Calls set_arena_queue from the correct thread so threading.local()
    stores the queue on the worker thread where emit_arena_event runs.
    """
    set_arena_queue(queue, loop)
    try:
        return run_arena(query=query, ticker=ticker, query_mode=query_mode)
    finally:
        clear_arena_queue()


@app.post("/memo/stream")
async def memo_stream(request: Request, memo_request: MemoRequest):
    """Stream Investment Committee memo via SSE.
    Runs the Arena debate, then post-processes the final state into a
    structured memo.  Emits arena progress events in real time, then
    emits arena_memo_ready when the structured memo is ready.

    Rate limits: 5 req/IP/hour, 50 req/day global (in-memory, single-worker only).
    Single Uvicorn worker required for in-memory counter accuracy.
    """
    ticker = memo_request.ticker.upper().strip()
    if not ticker or len(ticker) > 5 or not ticker.isalpha():
        raise HTTPException(status_code=400, detail="Invalid ticker symbol")

    # Pre-validate ticker against Financial Datasets API before starting Arena
    fd_key = os.getenv("FINANCIAL_DATASETS_API_KEY")
    if fd_key:
        try:
            fd_resp = requests.get(
                "https://api.financialdatasets.ai/financials",
                params={"ticker": ticker, "period": "annual", "limit": 1},
                headers={"X-API-KEY": fd_key},
                timeout=10,
            )
            resp_json = fd_resp.json()
            if fd_resp.status_code == 402 or "credits" in resp_json.get("message", "").lower() or "credits" in resp_json.get("error", "").lower():
                raise HTTPException(
                    status_code=503,
                    detail="Financial data service is temporarily unavailable. Please try again later.",
                )
            if fd_resp.status_code != 200 or not resp_json.get("financials"):
                raise HTTPException(
                    status_code=400,
                    detail=f"Unable to find financial data for {ticker}. Please verify the symbol.",
                )
        except HTTPException:
            raise
        except Exception:
            pass  # If validation call itself fails, proceed and let Arena handle it

    client_ip = request.client.host if request.client else "unknown"
    _check_memo_rate_limits(client_ip)

    async def generate():
        queue: asyncio.Queue = asyncio.Queue()
        loop = asyncio.get_running_loop()

        # Run Arena in thread pool — set_arena_queue called from within worker thread
        future = loop.run_in_executor(
            None,
            _run_arena_in_worker,
            f"Investment analysis for {ticker}",
            ticker,
            memo_request.query_mode,
            queue,
            loop,
        )

        # Drain SSE events from the arena while it runs; 180s timeout guard
        final_state = None
        try:
            arena_future = asyncio.wrap_future(future)
            while True:
                try:
                    event = await asyncio.wait_for(queue.get(), timeout=0.1)
                    if event is None:  # sentinel
                        break
                    yield f"data: {json.dumps(event)}\n\n"
                except asyncio.TimeoutError:
                    if future.done():
                        try:
                            final_state = future.result()
                        except Exception as exc:
                            yield f"data: {json.dumps({'type': 'error', 'error': str(exc)})}\n\n"
                            yield f"data: {json.dumps({'type': 'end'})}\n\n"
                            return
                        break
                    # Check overall 180s timeout
                    try:
                        await asyncio.wait_for(asyncio.shield(arena_future), timeout=0)
                    except asyncio.TimeoutError:
                        pass
        except asyncio.TimeoutError:
            yield f"data: {json.dumps({'type': 'error', 'error': 'Analysis timed out after 3 minutes — please try again'})}\n\n"
            yield f"data: {json.dumps({'type': 'end'})}\n\n"
            return

        if final_state is None:
            try:
                final_state = await asyncio.wait_for(
                    asyncio.wrap_future(future), timeout=180
                )
            except asyncio.TimeoutError:
                yield f"data: {json.dumps({'type': 'error', 'error': 'Analysis timed out after 3 minutes — please try again'})}\n\n"
                yield f"data: {json.dumps({'type': 'end'})}\n\n"
                return
            except Exception as exc:
                yield f"data: {json.dumps({'type': 'error', 'error': str(exc)})}\n\n"
                yield f"data: {json.dumps({'type': 'end'})}\n\n"
                return

        # Post-process: extract structured memo (Haiku call, ~3s)
        structured_memo = await loop.run_in_executor(None, extract_structured_memo, final_state)

        verdict, confidence = derive_verdict(
            final_state.get("consensus_score", 0.0),
            final_state.get("agent_signals", {}),
            final_state.get("next_action", ""),
        )

        yield f"data: {json.dumps({'type': 'arena_memo_ready', 'structured_memo': structured_memo, 'verdict': verdict, 'confidence': confidence, 'agent_signals': final_state.get('agent_signals', {}), 'debate_log': final_state.get('debate_log', [])})}\n\n"
        yield f"data: {json.dumps({'type': 'end'})}\n\n"

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        },
    )


class MemoSaveRequest(BaseModel):
    ticker: str
    verdict: str
    confidence: float
    structured_memo: dict
    checklist_answers: dict  # {why_now, exit_condition, max_position_size, quarterly_check_metric}


@app.post("/api/memo/save")
async def memo_save(payload: MemoSaveRequest, db: AsyncSession = Depends(get_db)):
    """Persist a completed memo with checklist answers. Returns share_slug."""
    required_keys = {"why_now", "exit_condition", "max_position_size", "quarterly_check_metric"}
    if not required_keys.issubset(payload.checklist_answers.keys()) or not all(
        str(v).strip() for v in payload.checklist_answers.values()
    ):
        raise HTTPException(status_code=422, detail="All 4 checklist fields are required to save")

    slug = secrets.token_urlsafe(6)
    content = json.dumps({
        "verdict": payload.verdict,
        "confidence": payload.confidence,
        "structured_memo": payload.structured_memo,
        "checklist_answers": payload.checklist_answers,
    })
    analysis = Analysis(
        ticker=payload.ticker.upper(),
        agent_type="memo",
        title=f"Investment Memo — {payload.ticker.upper()} ({payload.verdict})",
        content=content,
        tags=json.dumps(["memo"]),
        share_slug=slug,
        checklist_answers=json.dumps(payload.checklist_answers),
    )
    db.add(analysis)
    await db.commit()
    await db.refresh(analysis)
    return {"id": analysis.id, "share_slug": slug}


@app.get("/ticker/search")
async def ticker_search(q: str = ""):
    """Proxy Yahoo Finance autocomplete — returns matching tickers for a name or symbol."""
    q = q.strip()
    if not q or len(q) < 1:
        return []
    try:
        resp = requests.get(
            "https://query2.finance.yahoo.com/v1/finance/search",
            params={"q": q, "quotesCount": 7, "newsCount": 0, "enableFuzzyQuery": True},
            headers={"User-Agent": "Mozilla/5.0"},
            timeout=5,
        )
        data = resp.json()
        quotes = data.get("quotes", [])
        results = []
        for item in quotes:
            if item.get("quoteType") not in ("EQUITY", "ETF"):
                continue
            results.append({
                "symbol": item.get("symbol", ""),
                "name": item.get("longname") or item.get("shortname", ""),
                "exchange": item.get("exchDisp", ""),
                "type": item.get("quoteType", ""),
            })
        return results[:7]
    except Exception as e:
        logger.warning(f"Ticker search failed: {e}")
        return []


@app.get("/api/m/{slug}")
async def memo_by_slug(slug: str, db: AsyncSession = Depends(get_db)):
    """Public read-only endpoint for a shared memo URL."""
    from sqlalchemy import select as sa_select
    result = await db.execute(
        sa_select(Analysis).where(Analysis.share_slug == slug)
    )
    analysis = result.scalar_one_or_none()
    if analysis is None:
        raise HTTPException(status_code=404, detail="Memo not found")
    try:
        content = json.loads(analysis.content)
    except Exception:
        raise HTTPException(status_code=500, detail="Memo data corrupted")
    return {
        "ticker": analysis.ticker,
        "verdict": content.get("verdict"),
        "confidence": content.get("confidence"),
        "structured_memo": content.get("structured_memo"),
        "checklist_answers": content.get("checklist_answers"),
        "created_at": analysis.created_at.isoformat() if analysis.created_at else None,
    }


@app.get("/stock-chart/compare")
async def get_stock_chart_compare(tickers: str, period: str = "1M"):
    """
    Fetch stock chart data for multiple tickers for comparison.

    Args:
        tickers: Comma-separated ticker symbols (max 2), e.g. "AAPL,MSFT"
        period: Time period (1M, 6M, YTD, 1Y, 5Y, MAX)

    Returns:
        JSON with tickers list, quotes dict, and historical dict
    """
    try:
        fmp_key = os.getenv("FMP_API_KEY")
        if not fmp_key:
            raise HTTPException(status_code=500, detail="FMP_API_KEY not configured")

        ticker_list = [t.strip().upper() for t in tickers.split(",") if t.strip()]
        if len(ticker_list) == 0 or len(ticker_list) > 2:
            raise HTTPException(status_code=400, detail="Provide 1-2 comma-separated tickers")

        quote_url = "https://financialmodelingprep.com/stable/quote"
        hist_url = "https://financialmodelingprep.com/stable/historical-price-eod/full"

        async def load_compare_ticker(tkr: str) -> tuple[str, dict, List[Dict]]:
            quote_data_list, hist_result = await asyncio.gather(
                _fetch_json(quote_url, params={"symbol": tkr, "apikey": fmp_key}),
                _fetch_json(hist_url, params={"symbol": tkr, "apikey": fmp_key}),
                return_exceptions=True,
            )

            if isinstance(quote_data_list, Exception):
                logger.error(f"Failed to fetch compare chart data for {tkr}: {quote_data_list}")
                raise HTTPException(status_code=404, detail=f"Could not fetch data for ticker {tkr}")

            if not quote_data_list:
                raise HTTPException(status_code=404, detail=f"No quote data found for {tkr}")

            if isinstance(hist_result, Exception):
                logger.warning(f"Historical compare chart fetch failed for {tkr}: {hist_result}")
                hist_data = []
            else:
                hist_data = hist_result

            qd = quote_data_list[0]
            quote_payload = {
                "symbol": qd.get("symbol", tkr),
                "name": qd.get("name", ""),
                "exchange": qd.get("exchange", ""),
                "price": qd.get("price", 0),
                "changesPercentage": qd.get("changesPercentage", 0),
                "change": qd.get("change", 0),
                "dayHigh": qd.get("dayHigh", 0),
                "dayLow": qd.get("dayLow", 0),
                "volume": qd.get("volume", 0),
                "marketCap": qd.get("marketCap", 0),
                "open": qd.get("open", 0),
                "previousClose": qd.get("previousClose", 0),
                "yearHigh": qd.get("yearHigh", 0),
                "yearLow": qd.get("yearLow", 0),
                "avgVolume": qd.get("avgVolume", 0),
                "pe":   qd.get("pe", None),
                "eps":  qd.get("eps", None),
                "beta": qd.get("beta", None),
            }
            return tkr, quote_payload, filter_chart_data_by_period(hist_data, period)

        results = await asyncio.gather(*(load_compare_ticker(tkr) for tkr in ticker_list))
        quotes = {ticker: quote for ticker, quote, _ in results}
        historical = {ticker: history for ticker, _, history in results}

        return {
            "tickers": ticker_list,
            "quotes": quotes,
            "historical": historical,
        }

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


def filter_chart_data_by_period(data: Any, period: str) -> List[Dict]:
    """Filter historical chart data by time period"""
    if period == "1D":
        # Intraday data is already filtered by API (last trading day)
        return data if isinstance(data, list) else []

    # YTD uses Jan 1 of current year; all others use configured days mapping
    if period == "YTD":
        cutoff_date = datetime(datetime.now().year, 1, 1)
    else:
        days = CHART_PERIOD_DAYS.get(period, CHART_PERIOD_DAYS["1M"])
        cutoff_date = datetime.now() - timedelta(days=days)

    # Handle both list and dict responses
    historical = data.get("historical", data) if isinstance(data, dict) else data

    if not isinstance(historical, list):
        return []

    # Filter by date
    filtered = []
    for item in historical:
        try:
            # Parse date (format: "YYYY-MM-DD" or "YYYY-MM-DD HH:MM:SS")
            item_date_str = item.get("date", "").split(" ")[0]
            item_date = datetime.strptime(item_date_str, "%Y-%m-%d")

            if item_date >= cutoff_date:
                filtered.append(item)
        except (ValueError, AttributeError):
            continue

    return filtered


@app.get("/stock-chart/{ticker}")
async def get_stock_chart(ticker: str, period: str = "1M"):
    """
    Fetch stock chart data from FMP API

    Args:
        ticker: Stock ticker symbol (e.g., AAPL, MSFT)
        period: Time period (1D, 1W, 1M, 3M, 1Y, ALL)

    Returns:
        JSON with quote data and historical price data
    """
    try:
        fmp_key = os.getenv("FMP_API_KEY")
        if not fmp_key:
            raise HTTPException(status_code=500, detail="FMP_API_KEY not configured")

        ticker = ticker.upper()

        quote_url = "https://financialmodelingprep.com/stable/quote"
        if period == "1D":
            hist_url = "https://financialmodelingprep.com/stable/historical-chart/5min"
        else:
            hist_url = "https://financialmodelingprep.com/stable/historical-price-eod/full"

        quote_data_list, hist_result = await asyncio.gather(
            _fetch_json(quote_url, params={"symbol": ticker, "apikey": fmp_key}),
            _fetch_json(hist_url, params={"symbol": ticker, "apikey": fmp_key}),
            return_exceptions=True,
        )

        if isinstance(quote_data_list, Exception):
            print(f"[ERROR] Failed to fetch quote for {ticker}: {quote_data_list}")
            raise HTTPException(status_code=404, detail=f"Could not fetch data for ticker {ticker}")

        if not quote_data_list or len(quote_data_list) == 0:
            print(f"[ERROR] Empty quote response for {ticker}")
            raise HTTPException(status_code=404, detail=f"No quote data found for {ticker}")

        quote_data = quote_data_list[0]

        # Log the quote data structure for debugging
        print(f"[DEBUG] FMP Quote Response for {ticker}: {quote_data}")

        if isinstance(hist_result, Exception):
            print(f"[ERROR] Failed to fetch historical data for {ticker}: {hist_result}")
            hist_data = []
        else:
            hist_data = hist_result

        # Ensure required fields exist with fallbacks
        quote_data = {
            "symbol": quote_data.get("symbol", ticker),
            "name": quote_data.get("name", ""),
            "exchange": quote_data.get("exchange", ""),
            "price": quote_data.get("price", 0),
            "changesPercentage": quote_data.get("changesPercentage", 0),
            "change": quote_data.get("change", 0),
            "dayHigh": quote_data.get("dayHigh", 0),
            "dayLow": quote_data.get("dayLow", 0),
            "volume": quote_data.get("volume", 0),
            "marketCap": quote_data.get("marketCap", 0),
            "open": quote_data.get("open", 0),
            "previousClose": quote_data.get("previousClose", 0),
            "yearHigh": quote_data.get("yearHigh", 0),
            "yearLow": quote_data.get("yearLow", 0),
            "avgVolume": quote_data.get("avgVolume", 0),
            "pe":   quote_data.get("pe", None),
            "eps":  quote_data.get("eps", None),
            "beta": quote_data.get("beta", None),
        }

        # Log sample of historical data for debugging
        if isinstance(hist_data, list) and len(hist_data) > 0:
            print(f"[DEBUG] FMP Historical Response sample (first item): {hist_data[0]}")
        elif isinstance(hist_data, dict) and "historical" in hist_data:
            print(f"[DEBUG] FMP Historical Response sample (first item): {hist_data['historical'][0] if hist_data['historical'] else 'empty'}")

        # Filter historical data by period
        filtered_data = filter_chart_data_by_period(hist_data, period)

        return {
            "ticker": ticker,
            "quote": quote_data,
            "historical": filtered_data
        }

    except requests.exceptions.RequestException as e:
        raise HTTPException(status_code=502, detail=f"FMP API error: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/health")
async def health_check():
    """Detailed health check"""
    # Check if API keys are set
    api_keys = {
        "anthropic": bool(os.getenv("ANTHROPIC_API_KEY")),
        "financial_datasets": bool(os.getenv("FINANCIAL_DATASETS_API_KEY")),
        "tavily": bool(os.getenv("TAVILY_API_KEY")),
        "fred": bool(os.getenv("FRED_API_KEY")),
        "massive": bool(os.getenv("MASSIVE_API_KEY"))
    }

    return {
        "status": "healthy",
        "api_keys_configured": api_keys,
        "agents_cached": len(agents_cache),
        "timestamp": datetime.now().isoformat()
    }


ALLOWED_UPLOAD_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".csv"}
MAX_UPLOAD_SIZE = 10 * 1024 * 1024  # 10MB


def extract_text_from_file(filename: str, content: bytes) -> str:
    """Extract text from uploaded document based on file type."""
    import io
    ext = os.path.splitext(filename)[1].lower()

    if ext == ".pdf":
        from PyPDF2 import PdfReader
        reader = PdfReader(io.BytesIO(content))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(pages).strip()

    elif ext == ".docx":
        from docx import Document
        doc = Document(io.BytesIO(content))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())

    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        return "\n\n".join(t for t in texts if t.strip())

    elif ext == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        sheets_text = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append("\t".join(cells))
            if rows:
                sheets_text.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets_text)

    elif ext == ".csv":
        import csv
        reader = csv.reader(io.StringIO(content.decode("utf-8", errors="replace")))
        rows = ["\t".join(row) for row in reader]
        return "\n".join(rows)

    raise ValueError(f"Unsupported file type: {ext}")


@app.post("/upload-document")
async def upload_document(file: UploadFile = File(...)):
    """Upload a document and extract its text content."""
    if not file.filename:
        raise HTTPException(status_code=400, detail="No filename provided")

    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in ALLOWED_UPLOAD_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type: {ext}. Allowed: {', '.join(ALLOWED_UPLOAD_EXTENSIONS)}"
        )

    content = await file.read()
    if len(content) > MAX_UPLOAD_SIZE:
        raise HTTPException(status_code=413, detail="File too large. Maximum size is 10MB.")

    try:
        text = extract_text_from_file(file.filename, content)
    except Exception as e:
        logger.error(f"Failed to extract text from {file.filename}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to extract text: {str(e)}")

    return {"filename": file.filename, "content": text, "file_type": ext}


# =============================================================================
# REST Endpoints — Chat History (Sessions)
# =============================================================================

@app.get("/sessions")
async def list_sessions(limit: int = 50, db: AsyncSession = Depends(get_db)):
    """List sessions most-recent first."""
    result = await db.execute(
        select(DBSession).order_by(DBSession.last_active_at.desc()).limit(limit)
    )
    sessions = result.scalars().all()
    return [
        {
            "id": s.id,
            "title": s.title,
            "agent_type": s.agent_type,
            "created_at": s.created_at.isoformat(),
            "last_active_at": s.last_active_at.isoformat(),
        }
        for s in sessions
    ]


@app.get("/sessions/{session_id}")
async def get_session(session_id: str, db: AsyncSession = Depends(get_db)):
    """Get a session with all its messages."""
    result = await db.execute(select(DBSession).where(DBSession.id == session_id))
    session = result.scalar_one_or_none()
    if not session:
        raise HTTPException(status_code=404, detail="Session not found")

    msgs = await db.execute(
        select(DBMessage).where(DBMessage.session_id == session_id).order_by(DBMessage.created_at)
    )
    messages = msgs.scalars().all()
    return {
        "id": session.id,
        "title": session.title,
        "agent_type": session.agent_type,
        "created_at": session.created_at.isoformat(),
        "last_active_at": session.last_active_at.isoformat(),
        "messages": [
            {
                "id": m.id,
                "role": m.role,
                "content": m.content,
                "agent_type": m.agent_type,
                "ticker": m.ticker,
                "thinking_steps": json.loads(m.thinking_steps) if m.thinking_steps else [],
                "follow_ups": json.loads(m.follow_ups) if m.follow_ups else [],
                "chart_specs": json.loads(m.chart_specs) if m.chart_specs else {},
                "created_at": m.created_at.isoformat(),
            }
            for m in messages
        ],
    }


@app.delete("/sessions/{session_id}", status_code=204)
async def delete_session(session_id: str, db: AsyncSession = Depends(get_db)):
    """Delete a session and all its messages (cascade)."""
    result = await db.execute(select(DBSession).where(DBSession.id == session_id))
    session = result.scalar_one_or_none()
    if not session:
        raise HTTPException(status_code=404, detail="Session not found")
    await db.delete(session)
    await db.commit()
    # Evict any session-scoped agents from cache
    with agents_cache_lock:
        stale_keys = [k for k in agents_cache if k.endswith(f"_{session_id}")]
        for k in stale_keys:
            del agents_cache[k]
    return Response(status_code=204)


# =============================================================================
# REST Endpoints — Research Library (Analyses)
# =============================================================================

class AnalysisPatch(BaseModel):
    tags: Optional[list[str]] = None


@app.get("/analyses")
async def list_analyses(
    ticker: Optional[str] = None,
    tag: Optional[str] = None,
    q: Optional[str] = None,
    agent_type: Optional[str] = None,
    db: AsyncSession = Depends(get_db),
):
    """List saved analyses with optional filters."""
    stmt = select(Analysis).order_by(Analysis.created_at.desc())
    if ticker:
        stmt = stmt.where(Analysis.ticker == ticker.upper())
    if agent_type:
        stmt = stmt.where(Analysis.agent_type == agent_type)
    if q:
        stmt = stmt.where(
            or_(
                Analysis.title.ilike(f"%{q}%"),
                Analysis.content.ilike(f"%{q}%"),
            )
        )
    result = await db.execute(stmt)
    analyses = result.scalars().all()

    rows = []
    for a in analyses:
        tags = json.loads(a.tags) if a.tags else []
        if tag and tag not in tags:
            continue
        rows.append({
            "id": a.id,
            "ticker": a.ticker,
            "agent_type": a.agent_type,
            "title": a.title,
            "content_preview": a.content[:200],
            "tags": tags,
            "session_id": a.session_id,
            "created_at": a.created_at.isoformat(),
            "updated_at": a.updated_at.isoformat(),
        })
    return rows


@app.get("/analyses/{analysis_id}")
async def get_analysis(analysis_id: str, db: AsyncSession = Depends(get_db)):
    """Get a single analysis with full content."""
    result = await db.execute(select(Analysis).where(Analysis.id == analysis_id))
    a = result.scalar_one_or_none()
    if not a:
        raise HTTPException(status_code=404, detail="Analysis not found")
    return {
        "id": a.id,
        "ticker": a.ticker,
        "agent_type": a.agent_type,
        "title": a.title,
        "content": a.content,
        "tags": json.loads(a.tags) if a.tags else [],
        "session_id": a.session_id,
        "created_at": a.created_at.isoformat(),
        "updated_at": a.updated_at.isoformat(),
    }


@app.patch("/analyses/{analysis_id}")
async def update_analysis(analysis_id: str, patch: AnalysisPatch, db: AsyncSession = Depends(get_db)):
    """Update analysis tags."""
    result = await db.execute(select(Analysis).where(Analysis.id == analysis_id))
    a = result.scalar_one_or_none()
    if not a:
        raise HTTPException(status_code=404, detail="Analysis not found")
    if patch.tags is not None:
        a.tags = json.dumps(patch.tags)
        a.updated_at = datetime.now(timezone.utc)
    await db.commit()
    return {"id": a.id, "tags": json.loads(a.tags)}


@app.get("/analyses/{analysis_id}/export")
async def export_analysis(analysis_id: str, db: AsyncSession = Depends(get_db)):
    """Download analysis as a .md file."""
    result = await db.execute(select(Analysis).where(Analysis.id == analysis_id))
    a = result.scalar_one_or_none()
    if not a:
        raise HTTPException(status_code=404, detail="Analysis not found")
    date_str = a.created_at.strftime("%Y-%m-%d")
    filename = f"{a.ticker or 'analysis'}_{a.agent_type}_{date_str}.md"
    return Response(
        content=a.content,
        media_type="text/markdown",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@app.delete("/analyses/{analysis_id}", status_code=204)
async def delete_analysis(analysis_id: str, db: AsyncSession = Depends(get_db)):
    """Delete a saved analysis."""
    result = await db.execute(select(Analysis).where(Analysis.id == analysis_id))
    a = result.scalar_one_or_none()
    if not a:
        raise HTTPException(status_code=404, detail="Analysis not found")
    await db.delete(a)
    await db.commit()
    return Response(status_code=204)


# =============================================================================
# REST Endpoints — Watchlists
# =============================================================================

class WatchlistCreate(BaseModel):
    name: str = "My Watchlist"


class TickerAdd(BaseModel):
    ticker: str
    notes: Optional[str] = None


@app.post("/watchlists", status_code=201)
async def create_watchlist(body: WatchlistCreate, db: AsyncSession = Depends(get_db)):
    """Create a new watchlist."""
    wl = Watchlist(name=body.name)
    db.add(wl)
    await db.commit()
    return {"id": wl.id, "name": wl.name, "created_at": wl.created_at.isoformat()}


@app.get("/watchlists")
async def list_watchlists(db: AsyncSession = Depends(get_db)):
    """List all watchlists with their tickers."""
    result = await db.execute(select(Watchlist).order_by(Watchlist.created_at))
    watchlists = result.scalars().all()
    out = []
    for wl in watchlists:
        tickers_result = await db.execute(
            select(WatchlistTicker).where(WatchlistTicker.watchlist_id == wl.id).order_by(WatchlistTicker.added_at)
        )
        tickers = tickers_result.scalars().all()
        out.append({
            "id": wl.id,
            "name": wl.name,
            "created_at": wl.created_at.isoformat(),
            "tickers": [
                {"id": t.id, "ticker": t.ticker, "notes": t.notes, "added_at": t.added_at.isoformat()}
                for t in tickers
            ],
        })
    return out


@app.get("/watchlists/{watchlist_id}/tickers")
async def get_watchlist_tickers(watchlist_id: str, db: AsyncSession = Depends(get_db)):
    """Get tickers for a specific watchlist."""
    result = await db.execute(
        select(WatchlistTicker).where(WatchlistTicker.watchlist_id == watchlist_id).order_by(WatchlistTicker.added_at)
    )
    tickers = result.scalars().all()
    return [
        {"id": t.id, "ticker": t.ticker, "notes": t.notes, "added_at": t.added_at.isoformat()}
        for t in tickers
    ]


@app.post("/watchlists/{watchlist_id}/tickers", status_code=201)
async def add_ticker_to_watchlist(watchlist_id: str, body: TickerAdd, db: AsyncSession = Depends(get_db)):
    """Add a ticker to a watchlist."""
    result = await db.execute(select(Watchlist).where(Watchlist.id == watchlist_id))
    wl = result.scalar_one_or_none()
    if not wl:
        raise HTTPException(status_code=404, detail="Watchlist not found")

    # Check for duplicate
    existing = await db.execute(
        select(WatchlistTicker).where(
            WatchlistTicker.watchlist_id == watchlist_id,
            WatchlistTicker.ticker == body.ticker.upper(),
        )
    )
    if existing.scalar_one_or_none():
        raise HTTPException(status_code=409, detail=f"{body.ticker.upper()} already in watchlist")

    t = WatchlistTicker(watchlist_id=watchlist_id, ticker=body.ticker.upper(), notes=body.notes)
    db.add(t)
    await db.commit()
    return {"id": t.id, "ticker": t.ticker, "notes": t.notes, "added_at": t.added_at.isoformat()}


@app.delete("/watchlists/{watchlist_id}/tickers/{ticker}", status_code=204)
async def remove_ticker_from_watchlist(watchlist_id: str, ticker: str, db: AsyncSession = Depends(get_db)):
    """Remove a ticker from a watchlist."""
    result = await db.execute(
        select(WatchlistTicker).where(
            WatchlistTicker.watchlist_id == watchlist_id,
            WatchlistTicker.ticker == ticker.upper(),
        )
    )
    t = result.scalar_one_or_none()
    if not t:
        raise HTTPException(status_code=404, detail="Ticker not found in watchlist")
    await db.delete(t)
    await db.commit()
    return Response(status_code=204)


@app.delete("/watchlists/{watchlist_id}", status_code=204)
async def delete_watchlist(watchlist_id: str, db: AsyncSession = Depends(get_db)):
    """Delete a watchlist and all its tickers."""
    result = await db.execute(select(Watchlist).where(Watchlist.id == watchlist_id))
    wl = result.scalar_one_or_none()
    if not wl:
        raise HTTPException(status_code=404, detail="Watchlist not found")
    await db.delete(wl)
    await db.commit()
    return Response(status_code=204)


# =============================================================================
# REST Endpoints — Projects
# =============================================================================

class CreateProjectRequest(BaseModel):
    title: str
    thesis: str
    tickers: Optional[List[str]] = None


class ProjectMemoryPatch(BaseModel):
    memory_doc: str


class ProjectPatch(BaseModel):
    title: Optional[str] = None
    thesis: Optional[str] = None
    config: Optional[dict] = None
    status: Optional[str] = None


def _project_detail(p: Project, session_count: int = 0, document_count: int = 0) -> dict:
    config = json.loads(p.config) if p.config else {}
    return {
        "id": p.id,
        "title": p.title,
        "thesis": p.thesis,
        "config": config,
        "memory_doc": p.memory_doc,
        "status": p.status,
        "created_at": p.created_at.isoformat(),
        "updated_at": p.updated_at.isoformat(),
        "session_count": session_count,
        "document_count": document_count,
    }


def _project_summary(p: Project, session_count: int = 0, document_count: int = 0) -> dict:
    d = _project_detail(p, session_count, document_count)
    d.pop("memory_doc", None)
    return d


@app.post("/projects", status_code=201)
async def create_project(body: CreateProjectRequest, db: AsyncSession = Depends(get_db)):
    """Create a new investment thesis project."""
    from data.project_memory import initialize_memory_doc
    config_dict = normalize_project_config({"tickers": body.tickers, "preferred_agents": []})
    config = json.dumps(config_dict)
    memory_doc = initialize_memory_doc(body.title, body.thesis, tickers=config_dict.get("tickers", []))
    project = Project(
        title=body.title,
        thesis=body.thesis,
        config=config,
        memory_doc=memory_doc,
    )
    db.add(project)
    await db.commit()
    return _project_detail(project)


@app.get("/projects")
async def list_projects(db: AsyncSession = Depends(get_db)):
    """List active projects with session and document counts."""
    from sqlalchemy import func
    result = await db.execute(
        select(Project).where(Project.status == "active").order_by(Project.updated_at.desc())
    )
    projects = result.scalars().all()
    out = []
    for p in projects:
        sc_result = await db.execute(
            select(func.count()).select_from(ProjectSession).where(ProjectSession.project_id == p.id)
        )
        session_count = sc_result.scalar() or 0
        dc_result = await db.execute(
            select(func.count()).select_from(ProjectDocument).where(ProjectDocument.project_id == p.id)
        )
        document_count = dc_result.scalar() or 0
        out.append(_project_summary(p, session_count, document_count))
    return out


@app.get("/projects/{project_id}")
async def get_project(project_id: str, db: AsyncSession = Depends(get_db)):
    """Get full project detail including memory_doc."""
    from sqlalchemy import func
    result = await db.execute(select(Project).where(Project.id == project_id))
    p = result.scalar_one_or_none()
    if not p:
        raise HTTPException(status_code=404, detail="Project not found")
    sc_result = await db.execute(
        select(func.count()).select_from(ProjectSession).where(ProjectSession.project_id == p.id)
    )
    session_count = sc_result.scalar() or 0
    dc_result = await db.execute(
        select(func.count()).select_from(ProjectDocument).where(ProjectDocument.project_id == p.id)
    )
    document_count = dc_result.scalar() or 0
    return _project_detail(p, session_count, document_count)


@app.patch("/projects/{project_id}")
async def update_project(project_id: str, patch: ProjectPatch, db: AsyncSession = Depends(get_db)):
    """Update project title, thesis, config, or status."""
    from data.project_memory import sync_project_memory
    result = await db.execute(select(Project).where(Project.id == project_id))
    p = result.scalar_one_or_none()
    if not p:
        raise HTTPException(status_code=404, detail="Project not found")

    existing_config = json.loads(p.config) if p.config else {}
    next_title = patch.title if patch.title is not None else p.title
    next_thesis = patch.thesis if patch.thesis is not None else p.thesis
    next_config = existing_config

    if patch.title is not None:
        p.title = patch.title
    if patch.thesis is not None:
        p.thesis = patch.thesis
    if patch.config is not None:
        next_config = normalize_project_config(patch.config, existing=existing_config)
        p.config = json.dumps(next_config)
    if patch.status is not None:
        p.status = patch.status

    if patch.title is not None or patch.thesis is not None or patch.config is not None:
        p.memory_doc = sync_project_memory(
            p.memory_doc or "",
            title=next_title,
            thesis=next_thesis,
            tickers=next_config.get("tickers"),
        )
    p.updated_at = datetime.now(timezone.utc)
    await db.commit()
    return _project_detail(p)


@app.delete("/projects/{project_id}", status_code=204)
async def delete_project(project_id: str, db: AsyncSession = Depends(get_db)):
    """Archive a project and delete its Chroma collection."""
    result = await db.execute(select(Project).where(Project.id == project_id))
    p = result.scalar_one_or_none()
    if not p:
        raise HTTPException(status_code=404, detail="Project not found")
    p.status = "archived"
    p.updated_at = datetime.now(timezone.utc)
    await db.commit()
    # Best-effort Chroma cleanup
    try:
        from data.chroma_client import ProjectChromaClient
        chroma = ProjectChromaClient()
        await chroma.async_delete_collection(project_id)
    except Exception as e:
        logger.warning(f"Chroma cleanup failed for project {project_id}: {e}")
    return Response(status_code=204)


@app.get("/projects/{project_id}/memory")
async def get_project_memory(project_id: str, db: AsyncSession = Depends(get_db)):
    """Return the project memory document."""
    result = await db.execute(select(Project).where(Project.id == project_id))
    p = result.scalar_one_or_none()
    if not p:
        raise HTTPException(status_code=404, detail="Project not found")
    return {"memory_doc": p.memory_doc, "updated_at": p.updated_at.isoformat()}


@app.patch("/projects/{project_id}/memory")
async def patch_project_memory(project_id: str, body: ProjectMemoryPatch, db: AsyncSession = Depends(get_db)):
    """Manually overwrite the project memory document."""
    from data.project_memory import sync_project_memory

    result = await db.execute(select(Project).where(Project.id == project_id))
    p = result.scalar_one_or_none()
    if not p:
        raise HTTPException(status_code=404, detail="Project not found")
    now = datetime.now(timezone.utc)
    p.memory_doc = sync_project_memory(
        body.memory_doc,
        now_iso=now.replace(tzinfo=timezone.utc).isoformat(timespec="seconds"),
    )
    p.updated_at = now
    await db.commit()
    return {"memory_doc": p.memory_doc, "updated_at": p.updated_at.isoformat()}


@app.get("/projects/{project_id}/sessions")
async def list_project_sessions(project_id: str, db: AsyncSession = Depends(get_db)):
    """List sessions linked to a project."""
    result = await db.execute(select(Project).where(Project.id == project_id))
    if not result.scalar_one_or_none():
        raise HTTPException(status_code=404, detail="Project not found")
    ps_result = await db.execute(
        select(ProjectSession).where(ProjectSession.project_id == project_id).order_by(ProjectSession.created_at.desc())
    )
    project_sessions = ps_result.scalars().all()
    out = []
    for ps in project_sessions:
        s_result = await db.execute(select(DBSession).where(DBSession.id == ps.session_id))
        s = s_result.scalar_one_or_none()
        if s:
            out.append({
                "id": s.id,
                "title": s.title,
                "agent_type": s.agent_type,
                "created_at": s.created_at.isoformat(),
                "last_active_at": s.last_active_at.isoformat(),
            })
    return out


# =============================================================================
# REST Endpoints — Project Documents
# =============================================================================

@app.post("/projects/{project_id}/documents", status_code=201)
async def upload_project_document(
    project_id: str,
    file: UploadFile = File(...),
    db: AsyncSession = Depends(get_db),
):
    """Upload a document to a project: extract text, embed in Chroma, save record."""
    result = await db.execute(select(Project).where(Project.id == project_id))
    p = result.scalar_one_or_none()
    if not p:
        raise HTTPException(status_code=404, detail="Project not found")

    if not file.filename:
        raise HTTPException(status_code=400, detail="No filename provided")

    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in ALLOWED_UPLOAD_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type: {ext}. Allowed: {', '.join(ALLOWED_UPLOAD_EXTENSIONS)}"
        )

    content = await file.read()
    if len(content) > MAX_UPLOAD_SIZE:
        raise HTTPException(status_code=413, detail="File too large. Maximum size is 10MB.")

    try:
        raw_text = extract_text_from_file(file.filename, content)
    except Exception as e:
        logger.error(f"Failed to extract text from {file.filename}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to extract text: {str(e)}")

    # Create document record to get an ID
    import uuid as _uuid_mod
    doc_id = str(_uuid_mod.uuid4())
    doc = ProjectDocument(
        id=doc_id,
        project_id=project_id,
        filename=file.filename,
        file_type=ext,
        raw_text=raw_text,
        chunk_count=0,
        chroma_ids="[]",
    )
    db.add(doc)
    await db.flush()

    # Embed chunks in Chroma
    chroma_ids: List[str] = []
    try:
        from data.chroma_client import ProjectChromaClient
        chroma = ProjectChromaClient()
        chroma_ids = await chroma.async_add_document_chunks(project_id, doc_id, file.filename, raw_text)
        doc.chunk_count = len(chroma_ids)
        doc.chroma_ids = json.dumps(chroma_ids)
    except Exception as e:
        logger.warning(f"Chroma embedding failed for {file.filename}: {e}")

    # Generate summary and patch memory_doc
    try:
        from langchain_anthropic import ChatAnthropic
        from data.project_memory import (
            format_document_summary_entry,
            generate_document_summary,
            patch_memory_section,
        )
        llm = ChatAnthropic(model="claude-haiku-4-5-20251001", max_tokens=400)
        summary = generate_document_summary(file.filename, raw_text, llm)
        p.memory_doc = patch_memory_section(
            p.memory_doc or "",
            "Uploaded Document Summaries",
            format_document_summary_entry(file.filename, summary, document_id=doc_id),
            mode="append",
        )
        p.updated_at = datetime.now(timezone.utc)
    except Exception as e:
        logger.warning(f"Document summary generation failed for {file.filename}: {e}")

    await db.commit()
    await db.refresh(doc)

    return {
        "id": doc.id,
        "project_id": doc.project_id,
        "filename": doc.filename,
        "file_type": doc.file_type,
        "chunk_count": doc.chunk_count,
        "uploaded_at": doc.uploaded_at.isoformat(),
    }


@app.get("/projects/{project_id}/documents")
async def list_project_documents(project_id: str, db: AsyncSession = Depends(get_db)):
    """List documents for a project (no raw_text in response)."""
    result = await db.execute(select(Project).where(Project.id == project_id))
    if not result.scalar_one_or_none():
        raise HTTPException(status_code=404, detail="Project not found")

    docs_result = await db.execute(
        select(ProjectDocument).where(ProjectDocument.project_id == project_id)
        .order_by(ProjectDocument.uploaded_at.desc())
    )
    docs = docs_result.scalars().all()
    return [
        {
            "id": d.id,
            "project_id": d.project_id,
            "filename": d.filename,
            "file_type": d.file_type,
            "chunk_count": d.chunk_count,
            "uploaded_at": d.uploaded_at.isoformat(),
        }
        for d in docs
    ]


@app.delete("/projects/{project_id}/documents/{doc_id}")
async def delete_project_document(project_id: str, doc_id: str, db: AsyncSession = Depends(get_db)):
    """Delete a project document and its Chroma chunks."""
    from data.project_memory import remove_document_summary

    result = await db.execute(
        select(ProjectDocument).where(
            ProjectDocument.id == doc_id,
            ProjectDocument.project_id == project_id,
        )
    )
    doc = result.scalar_one_or_none()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")

    # Delete Chroma chunks
    if doc.chroma_ids:
        try:
            from data.chroma_client import ProjectChromaClient
            chroma = ProjectChromaClient()
            ids = json.loads(doc.chroma_ids)
            if ids:
                await chroma.async_delete_chunks(project_id, ids)
        except Exception as e:
            logger.warning(f"Chroma chunk deletion failed for doc {doc_id}: {e}")

    project_result = await db.execute(select(Project).where(Project.id == project_id))
    project = project_result.scalar_one_or_none()
    if project is not None:
        project.memory_doc = remove_document_summary(
            project.memory_doc or "",
            doc.filename,
            document_id=doc.id,
        )
        project.updated_at = datetime.now(timezone.utc)

    await db.delete(doc)
    await db.commit()
    return Response(status_code=204)


# ---------------------------------------------------------------------------
# Static file serving for production (serves React build)
# ---------------------------------------------------------------------------

# Mount static files if build exists (production mode)
if os.path.isdir(FRONTEND_BUILD_DIR):
    # Serve static assets (JS, CSS, images)
    app.mount("/assets", StaticFiles(directory=os.path.join(FRONTEND_BUILD_DIR, "assets")), name="static")
    
    # Catch-all route for SPA - must be registered last
    @app.get("/{full_path:path}", include_in_schema=False)
    async def serve_spa(full_path: str):
        """Serve React SPA for all non-API routes"""
        # Don't serve SPA for API routes
        if full_path.startswith("api/") or full_path in ["docs", "redoc", "openapi.json"]:
            raise HTTPException(status_code=404)
        
        # Check if requesting a specific file that exists
        file_path = os.path.join(FRONTEND_BUILD_DIR, full_path)
        if os.path.isfile(file_path):
            return FileResponse(file_path)
        
        # Otherwise serve index.html for SPA routing
        index_path = os.path.join(FRONTEND_BUILD_DIR, "index.html")
        if os.path.isfile(index_path):
            return FileResponse(index_path)
        
        raise HTTPException(status_code=404, detail="Not found")
    
    logger.info(f"Serving frontend from {FRONTEND_BUILD_DIR}")
else:
    logger.info("Frontend build not found - API-only mode")


if __name__ == "__main__":
    import uvicorn

    # Check for required API keys
    if not os.getenv("ANTHROPIC_API_KEY"):
        print("ERROR: ANTHROPIC_API_KEY not set in .env file")
        sys.exit(1)

    print("Starting Financial Analysis API Server...")
    print("API will be available at: http://localhost:8000")
    print("API documentation: http://localhost:8000/docs")
    print("Available agents: Equity Analyst, Finance Q&A, Market Analyst, Portfolio Analyzer")

    uvicorn.run(
        "api_server:app",
        host="0.0.0.0",
        port=8000,
        reload=True,
        log_level="info"
    )
